"""
render.py — PPT 渲染器（仅依赖 python-pptx）

用法:
    python render.py outline.json [template.pptx] [output.pptx]

输入 outline.json 结构（这就是大模型需要产出的"结构化大纲"）:
{
  "slides": [
    {"type": "cover",   "title": "...", "subtitle": "...", "date": "..."},
    {"type": "bullets", "title": "...", "bullets": ["...", "...", "..."]},
    {"type": "two_col", "title": "...", "body": "...", "img:main": "/path/or/url.png"},
    ...
  ]
}

工作原理:
  - 模板里每种页型是一张范例页，页型 id 写在该页 notes 里（TYPE=xxx）。
  - 对大纲里的每一页，复制对应页型的范例页，按 token 填入内容，追加到末尾。
  - 全部填完后，删掉开头那些"范例页"，剩下的就是成品。
样式完全由范例页保真带过来，代码不碰任何排版。
"""
import sys, json, copy, re, io, os, urllib.request
from pptx import Presentation
from pptx.parts.image import Image as PptxImage

NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
}
TOKEN = re.compile(r"\{\{([^}]+)\}\}")
IMG_REL = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image"


# ---------- 识别页型 ----------
def slide_type(slide):
    if not slide.has_notes_slide:
        return None
    txt = slide.notes_slide.notes_text_frame.text.strip()
    m = re.search(r"TYPE\s*=\s*(\S+)", txt)
    return m.group(1) if m else None


# ---------- 复制一整页（含图片关系修复） ----------
def copy_slide(prs, src):
    new = prs.slides.add_slide(src.slide_layout)
    # 清掉 add_slide 自动带出的版式占位符，避免和 deepcopy 的形状重叠
    for ph in list(new.placeholders):
        ph._element.getparent().remove(ph._element)
    for shp in src.shapes:
        el = copy.deepcopy(shp._element)
        new.shapes._spTree.append(el)
        # 修复图片：deepcopy 出的 r:embed 仍指向源页的 rId，需在新页重建关系
        for blip in el.findall(".//a:blip", NS):
            rid = blip.get(f"{{{NS['r']}}}embed")
            if not rid:
                continue
            img_part = src.part.related_part(rid)
            new_rid = new.part.relate_to(img_part, IMG_REL)
            blip.set(f"{{{NS['r']}}}embed", new_rid)
    # 复制背景填充
    src_bg = src._element.find("{http://schemas.openxmlformats.org/presentationml/2006/main}cSld/"
                               "{http://schemas.openxmlformats.org/presentationml/2006/main}bg")
    if src_bg is not None:
        cSld = new._element.find("{http://schemas.openxmlformats.org/presentationml/2006/main}cSld")
        cSld.insert(0, copy.deepcopy(src_bg))
    return new


def delete_slide(prs, slide):
    sldId_lst = prs.slides._sldIdLst
    rId = slide.part.partname  # not used directly; remove by matching
    for sldId in list(sldId_lst):
        if prs.part.related_part(sldId.get(f"{{{NS['r']}}}id")) is slide.part:
            prs.part.drop_rel(sldId.get(f"{{{NS['r']}}}id"))
            sldId_lst.remove(sldId)
            break


# ---------- 文本替换（跨 run 安全） ----------
def _replace_in_paragraph(p, mapping):
    """把段落里所有 {{key}} 替换；合并到首 run，保留首 run 格式。"""
    runs = p.runs
    if not runs:
        return
    full = "".join(r.text for r in runs)
    if "{{" not in full:
        return
    new_text = TOKEN.sub(lambda m: str(mapping.get(m.group(1).strip(), m.group(0))), full)
    runs[0].text = new_text
    for r in runs[1:]:
        r.text = ""


def _iter_text_frames(shapes):
    for shp in shapes:
        if shp.has_text_frame:
            yield shp.text_frame
        if shp.shape_type == 6:  # GROUP
            yield from _iter_text_frames(shp.shapes)
        if shp.has_table:
            for row in shp.table.rows:
                for cell in row.cells:
                    yield cell.text_frame


def _expand_list_paragraph(tf, key, items):
    """某段落整段=={{*key}} 时，按 items 展开成多段，克隆该段格式。"""
    target = None
    for p in tf.paragraphs:
        if "".join(r.text for r in p.runs).strip() == "{{*%s}}" % key:
            target = p
            break
    if target is None:
        return False
    tmpl_pPr = target._p.find("{%s}pPr" % NS["a"])
    tmpl_run = target.runs[0]._r if target.runs else None
    parent = target._p.getparent()
    idx = list(parent).index(target._p)
    parent.remove(target._p)
    from pptx.oxml.ns import qn
    for i, item in enumerate(items):
        np = copy.deepcopy(target._p)
        # 清空旧 run 文本，写入新值
        np_runs = np.findall(qn("a:r"))
        if np_runs:
            t = np_runs[0].find(qn("a:t"))
            if t is not None:
                t.text = str(item)
            for extra in np_runs[1:]:
                np.remove(extra)
        parent.insert(idx + i, np)
    return True


def fill_slide(slide, data):
    # 1) 先处理列表占位（{{*key}}）
    for tf in _iter_text_frames(slide.shapes):
        for key, val in data.items():
            if isinstance(val, list):
                _expand_list_paragraph(tf, key, val)
    # 2) 标量占位（{{key}}）
    scalar = {k: v for k, v in data.items() if not isinstance(v, (list, dict))
              and not k.startswith("img:")}
    for tf in _iter_text_frames(slide.shapes):
        for p in tf.paragraphs:
            _replace_in_paragraph(p, scalar)
    # 3) 图片占位（形状名 == {{img:key}}）
    for shp in list(slide.shapes):
        m = re.fullmatch(r"\{\{(img:[^}]+)\}\}", shp.name or "")
        if m and m.group(1) in data:
            _swap_image(slide, shp, data[m.group(1)])


def _load_bytes(src):
    if str(src).startswith(("http://", "https://")):
        return urllib.request.urlopen(src, timeout=20).read()
    with open(src, "rb") as f:
        return f.read()


def _swap_image(slide, pic, src):
    """保留原图片框的位置/尺寸，仅替换图像数据。"""
    blob = _load_bytes(src)
    image_part, rId = slide.part.get_or_add_image_part(io.BytesIO(blob))
    blip = pic._element.find(".//a:blip", NS)
    old = blip.get(f"{{{NS['r']}}}embed")
    blip.set(f"{{{NS['r']}}}embed", rId)
    # 形状名清掉 token，避免成品里残留标记
    pic.name = pic.name.replace("{{", "").replace("}}", "")


# ---------- 主流程 ----------
def render(outline_path, template="template.pptx", output="output.pptx"):
    outline = json.load(open(outline_path, encoding="utf-8"))
    prs = Presentation(template)

    # 建立 页型 -> 范例页 索引；记录原始范例页以便最后删除
    exemplars = {}
    originals = list(prs.slides)
    for s in originals:
        t = slide_type(s)
        if t and t not in exemplars:
            exemplars[t] = s

    missing = {s.get("type") for s in outline["slides"]} - set(exemplars)
    if missing:
        raise SystemExit(f"模板缺少这些页型的范例页: {sorted(missing)}")

    for spec in outline["slides"]:
        new = copy_slide(prs, exemplars[spec["type"]])
        fill_slide(new, {k: v for k, v in spec.items() if k != "type"})

    for s in originals:           # 删掉开头的范例页
        delete_slide(prs, s)

    prs.save(output)
    print(f"saved {output}  ({len(outline['slides'])} 页)")


if __name__ == "__main__":
    args = sys.argv[1:]
    outline = args[0] if args else "example_outline.json"
    template = args[1] if len(args) > 1 else "template.pptx"
    output = args[2] if len(args) > 2 else "output.pptx"
    render(outline, template, output)
